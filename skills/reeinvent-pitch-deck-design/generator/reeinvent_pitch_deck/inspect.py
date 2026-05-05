"""Inspect / edit a cloned master deck. Path B workflow.

Two operations Claude needs to do reliably when working with a clone of the
Reeinvent master deck:

1. `inspect`: list every slide with its index, layout name, and the text
   content of each shape. Used to plan edits before touching anything.

2. `replace_text`: substitute strings inside text frames without altering
   layout, fonts, colors, gradients, or images. Operates per-slide or
   globally. Run order is important: replace from longest pattern to
   shortest so substrings don't shadow longer matches.

Both operations preserve every brand element on the master because they
only mutate text run content; shape geometry and styling stay untouched.
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation


def _iter_text_frames(shape):
    """Yield text_frame objects from a shape, recursing into groups."""
    if shape.has_text_frame:
        yield shape.text_frame
    if shape.shape_type == 6:  # group
        for sub in shape.shapes:
            yield from _iter_text_frames(sub)


def inspect_deck(path: Path) -> list[dict]:
    """Return a list of slide dicts: {index, layout, shapes:[{name, text}]}."""
    prs = Presentation(str(path))
    out = []
    for idx, slide in enumerate(prs.slides, start=1):
        layout = getattr(slide.slide_layout, "name", "?") or "?"
        shapes_info = []
        for shape in slide.shapes:
            text_pieces = []
            for tf in _iter_text_frames(shape):
                if tf.text.strip():
                    text_pieces.append(tf.text)
            if text_pieces:
                shapes_info.append({
                    "name": shape.name,
                    "text": " | ".join(text_pieces),
                })
        out.append({"index": idx, "layout": layout, "shapes": shapes_info})
    return out


def replace_text_in_deck(
    path: Path,
    output: Path,
    replacements: dict[str, str],
    only_slides: set[int] | None = None,
) -> int:
    """Run text replacements across the deck. Returns count of replacements made.

    `replacements` is {old_string: new_string}. Replacements run from longest
    `old_string` to shortest so a longer pattern is not eaten by a shorter one.

    `only_slides` (1-based indices) restricts the operation; None means all.

    Replacement is at the RUN level: we iterate runs and only replace within
    a single run when the entire `old_string` is contained in that run's
    text. This preserves typography (font, weight, color) on the kept
    portion of the run. If a pattern spans multiple runs (because the master
    deck uses inline bold), it will not match a single run; in that case the
    function joins the paragraph's runs as a single string for matching, and
    when found, the replacement goes into the FIRST run in the paragraph and
    later runs in the paragraph that participated in the match are cleared.
    This keeps the paragraph's leading typography while replacing the text.
    """
    prs = Presentation(str(path))
    count = 0
    sorted_keys = sorted(replacements.keys(), key=len, reverse=True)

    for sl_idx, slide in enumerate(prs.slides, start=1):
        if only_slides is not None and sl_idx not in only_slides:
            continue
        for shape in slide.shapes:
            for tf in _iter_text_frames(shape):
                for paragraph in tf.paragraphs:
                    for old in sorted_keys:
                        new = replacements[old]
                        # Try run-local replacement first (preserves run-level styling).
                        for run in paragraph.runs:
                            if old in run.text:
                                run.text = run.text.replace(old, new)
                                count += 1
                        # If old still appears across run boundaries, fall back
                        # to paragraph-level replacement.
                        full = "".join(r.text for r in paragraph.runs)
                        if old in full:
                            replaced = full.replace(old, new)
                            runs = list(paragraph.runs)
                            if runs:
                                runs[0].text = replaced
                                for r in runs[1:]:
                                    r.text = ""
                                count += 1

    prs.save(str(output))
    return count


def _cmd_inspect(args) -> int:
    deck = inspect_deck(Path(args.deck))
    if args.json:
        print(json.dumps(deck, indent=2, ensure_ascii=False))
    else:
        for s in deck:
            print(f"=== Slide {s['index']} (layout: {s['layout']}) ===")
            for shape in s["shapes"]:
                snippet = shape["text"][:120]
                print(f"  [{shape['name']}] {snippet}")
            print()
    return 0


def _cmd_replace(args) -> int:
    raw = json.loads(Path(args.replacements_json).read_text())
    if not isinstance(raw, dict):
        print("replace: replacements file must be a JSON object {old: new}", file=sys.stderr)
        return 1
    only = None
    if args.slides:
        only = {int(x) for x in args.slides.split(",")}
    out = Path(args.output) if args.output else Path(args.deck)
    count = replace_text_in_deck(Path(args.deck), out, raw, only_slides=only)
    print(f"replace: {count} replacements made -> {out}")
    return 0


def main(argv: list[str] | None = None) -> int:
    p = argparse.ArgumentParser(prog="reeinvent-deck-template")
    sub = p.add_subparsers(dest="command", required=True)

    pi = sub.add_parser("inspect", help="List slides + text for planning edits")
    pi.add_argument("deck")
    pi.add_argument("--json", action="store_true", help="Output JSON instead of human-readable")
    pi.set_defaults(func=_cmd_inspect)

    pr = sub.add_parser("replace", help="Apply text replacements across the deck")
    pr.add_argument("deck")
    pr.add_argument("replacements_json", help="JSON file mapping old->new strings")
    pr.add_argument("-o", "--output", help="output deck (default: in place)")
    pr.add_argument("--slides", help="comma-separated 1-based slide indices to limit to (default: all)")
    pr.set_defaults(func=_cmd_replace)

    args = p.parse_args(argv)
    return args.func(args)


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
