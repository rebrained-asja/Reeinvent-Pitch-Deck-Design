"""Build a Reeinvent-branded Presentation: 16:9 canvas, brand theme colors,
Roboto theme fonts. Returns a python-pptx Presentation ready for builders to
add slides into.

Customizing python-pptx's default master is done by patching the master's
theme1.xml after creation: replace the color scheme (clrScheme) and the font
scheme (fontScheme) with brand-correct values. Existing default layouts stay
in place so PowerPoint's "New Slide" menu still works for the presenter.
"""

from __future__ import annotations

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

from reeinvent_pitch_deck.theme import (
    CORE_BLUE,
    CORE_VIOLET,
    DEEP_NAVY,
    INK,
    MID_BLUE,
    MID_VIOLET,
    OFF_WHITE,
    SKY_BLUE,
    SLIDE_HEIGHT_IN,
    SLIDE_WIDTH_IN,
    SOFT_VIOLET,
    WHITE,
)


A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# Maps PowerPoint theme color slot -> brand hex. Order matches DESIGN.md
# section 12.1 rule 4 ("Core Blue mapped to ACCENT_1").
COLOR_SCHEME = [
    ("dk1", INK),
    ("lt1", OFF_WHITE),
    ("dk2", DEEP_NAVY),
    ("lt2", WHITE),
    ("accent1", CORE_BLUE),
    ("accent2", MID_BLUE),
    ("accent3", SKY_BLUE),
    ("accent4", CORE_VIOLET),
    ("accent5", MID_VIOLET),
    ("accent6", SOFT_VIOLET),
    ("hlink", CORE_BLUE),
    ("folHlink", MID_VIOLET),
]


def _replace_color_scheme(theme_root: etree._Element) -> None:
    theme_elements = theme_root.find(qn("a:themeElements"))
    if theme_elements is None:
        raise RuntimeError("theme1.xml has no themeElements; unsupported PPTX shape")
    old = theme_elements.find(qn("a:clrScheme"))
    if old is not None:
        theme_elements.remove(old)
    new = etree.SubElement(
        theme_elements,
        qn("a:clrScheme"),
        {"name": "Reeinvent"},
    )
    # System colors (dk1, lt1) need <a:sysClr> wrappers per OOXML; brand colors
    # use <a:srgbClr>. PowerPoint accepts srgbClr for all 12 slots in practice,
    # which keeps this simple and theme-edit safe.
    for slot, hex_val in COLOR_SCHEME:
        wrapper = etree.SubElement(new, qn(f"a:{slot}"))
        etree.SubElement(wrapper, qn("a:srgbClr"), {"val": hex_val})
    # Move clrScheme to the front of themeElements so it precedes fontScheme
    # and fmtScheme (required ordering).
    theme_elements.remove(new)
    theme_elements.insert(0, new)


def _replace_font_scheme(theme_root: etree._Element) -> None:
    theme_elements = theme_root.find(qn("a:themeElements"))
    if theme_elements is None:
        return
    old = theme_elements.find(qn("a:fontScheme"))
    if old is not None:
        theme_elements.remove(old)
    fs = etree.SubElement(theme_elements, qn("a:fontScheme"), {"name": "Reeinvent"})
    for major_or_minor in ("majorFont", "minorFont"):
        block = etree.SubElement(fs, qn(f"a:{major_or_minor}"))
        etree.SubElement(block, qn("a:latin"), {"typeface": "Roboto"})
        etree.SubElement(block, qn("a:ea"), {"typeface": ""})
        etree.SubElement(block, qn("a:cs"), {"typeface": ""})
    # Re-order: clrScheme, fontScheme, fmtScheme (existing).
    fmt_scheme = theme_elements.find(qn("a:fmtScheme"))
    children = list(theme_elements)
    for c in children:
        theme_elements.remove(c)
    cs = next((c for c in children if c.tag == qn("a:clrScheme")), None)
    if cs is not None:
        theme_elements.append(cs)
    theme_elements.append(fs)
    if fmt_scheme is not None:
        theme_elements.append(fmt_scheme)


def _patch_master_theme(prs: Presentation) -> None:
    """Apply Reeinvent color + font scheme to EVERY theme part in the package.

    A presentation has multiple theme parts: one for the main slide master
    (theme1.xml, a base Part with `_blob` as source of truth), one for the
    notes master (theme2.xml, an XmlPart with `_element` as source of truth),
    plus the handout master if present. The two part classes differ in how
    you mutate them, so this handles both.
    """
    from pptx.opc.package import XmlPart  # imported lazily to avoid cycles

    seen: set[str] = set()
    for part in prs.part.package.iter_parts():
        partname = str(getattr(part, "partname", ""))
        if not partname.startswith("/ppt/theme/") or not partname.endswith(".xml"):
            continue
        if partname in seen:
            continue
        seen.add(partname)

        # Source of truth differs by part class.
        if isinstance(part, XmlPart):
            theme_root = part._element
            _replace_color_scheme(theme_root)
            _replace_font_scheme(theme_root)
            # XmlPart serializes _element on read; mutating in place is enough.
        else:
            theme_root = etree.fromstring(part.blob)
            _replace_color_scheme(theme_root)
            _replace_font_scheme(theme_root)
            part._blob = etree.tostring(
                theme_root, xml_declaration=True, encoding="UTF-8"
            )


def new_branded_presentation() -> Presentation:
    """Return a fresh Presentation with the canvas set to 16:9 widescreen.

    Theme patching happens later via finalize_branded_presentation() because
    the notes master + its theme are created lazily when speaker notes are
    first added; patching upfront would miss them.
    """
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    return prs


def finalize_branded_presentation(prs: Presentation) -> None:
    """Apply Reeinvent theme to every theme part in the package. Call this
    AFTER all slides have been added so any lazy-created masters (notes,
    handout) get their theme patched too.
    """
    _patch_master_theme(prs)


BLANK_LAYOUT_INDEX = 6  # python-pptx default master ships layouts 0..10; index 6 is "Blank"


def add_blank_slide(prs: Presentation):
    """Add a slide using the master's blank layout. Builders draw everything
    explicitly on top, so the layout starts clean.
    """
    layout = prs.slide_layouts[BLANK_LAYOUT_INDEX]
    return prs.slides.add_slide(layout)
