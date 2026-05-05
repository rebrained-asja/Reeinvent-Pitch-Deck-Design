"""Drawing primitives shared across all archetype builders.

Every brand-controlled element (gradient fills, gradient text runs, gradient
stripes, brand-stamp logo, arrow watermark, gradient pill, chat bubble, arrow
bullet marker) is built here so changes propagate. Builders compose these.
"""

from __future__ import annotations

from typing import Iterable

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.shapes.autoshape import Shape
from pptx.slide import Slide
from pptx.util import Emu, Inches, Pt

from reeinvent_pitch_deck import assets
from reeinvent_pitch_deck.theme import (
    CORE_BLUE,
    Gradient,
    INK,
    LOGO_STAMP_HEIGHT_IN,
    LOGO_STAMP_RIGHT_IN,
    LOGO_STAMP_TOP_IN,
    OFF_WHITE,
    SIGNATURE_GRADIENT,
    SLIDE_HEIGHT_IN,
    SLIDE_WIDTH_IN,
    TYPE_SCALE,
    TypeRole,
    WHITE,
    is_bold_weight,
    weight_to_typeface,
)


A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
NSMAP = {"a": A_NS}


def _qn(tag: str) -> str:
    return qn(tag)


def _hex_to_rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str)


def _make_grad_fill_xml(grad: Gradient, opacity_pct: float | None = None) -> etree._Element:
    """Build an <a:gradFill> element for shape OR text fills.

    Both at 30 deg always. Two stops: 0% and 100%.

    OOXML angles use 60000ths of a degree. 30 deg = 1800000.
    """
    angle_60k = grad.angle_deg * 60000
    opacity_attrs = ""
    if opacity_pct is not None:
        # alpha is also in thousandths of percent (100% = 100000)
        alpha = int(opacity_pct * 1000)
        opacity_attrs = f'<a:alpha val="{alpha}"/>'
    xml = (
        '<a:gradFill xmlns:a="{ns}" rotWithShape="1">'
        '<a:gsLst>'
        '<a:gs pos="0"><a:srgbClr val="{c0}">{op}</a:srgbClr></a:gs>'
        '<a:gs pos="100000"><a:srgbClr val="{c1}">{op}</a:srgbClr></a:gs>'
        '</a:gsLst>'
        '<a:lin ang="{ang}" scaled="1"/>'
        '</a:gradFill>'
    ).format(ns=A_NS, c0=grad.stop_0, c1=grad.stop_100, op=opacity_attrs, ang=angle_60k)
    return etree.fromstring(xml)


def _replace_fill(parent: etree._Element, new_fill: etree._Element) -> None:
    """Replace any solidFill/gradFill/noFill/blipFill/pattFill child with new_fill."""
    fill_tags = {qn(f"a:{t}") for t in ("solidFill", "gradFill", "noFill", "blipFill", "pattFill")}
    for child in list(parent):
        if child.tag in fill_tags:
            parent.remove(child)
    parent.insert(0, new_fill)


def fill_shape_with_gradient(shape: Shape, grad: Gradient) -> None:
    """Apply a 30 deg native PowerPoint gradient fill to a shape."""
    sppr = shape.fill._xPr.find(qn("a:spPr"))  # type: ignore[attr-defined]
    if sppr is None:
        sppr = shape.fill._xPr  # text-frame fills don't wrap in spPr
    _replace_fill(sppr, _make_grad_fill_xml(grad))


def fill_run_with_gradient(run, grad: Gradient) -> None:
    """Apply a 30 deg native PowerPoint gradient fill to a text run."""
    rpr = run._r.get_or_add_rPr()
    _replace_fill(rpr, _make_grad_fill_xml(grad))


def set_run_typography(
    run,
    text: str,
    role: TypeRole,
    color_hex: str | None = None,
    italic: bool | None = None,
    weight_override: int | None = None,
    typeface_override: str | None = None,
) -> None:
    """Apply Reeinvent type rules to a run.

    Sets the text, font name (Roboto family), size, weight, italic, and color.
    Color is solid only here; for gradient text, call fill_run_with_gradient
    afterwards (it overrides the solid fill).
    """
    weight = weight_override if weight_override is not None else role.weight
    typeface = typeface_override if typeface_override is not None else weight_to_typeface(weight)
    run.text = text
    font = run.font
    font.name = typeface
    font.size = Pt(role.size_pt)
    font.bold = is_bold_weight(weight)
    font.italic = role.italic if italic is None else italic
    if color_hex is not None:
        font.color.rgb = _hex_to_rgb(color_hex)
    # Force the typeface as the LATIN typeface so PowerPoint maps to it instead
    # of the slide-master fallback. python-pptx's font.name does this for the
    # latin typeface element but only when the run already has rPr; ensure it.
    rpr = run._r.get_or_add_rPr()
    latin = rpr.find(qn("a:latin"))
    if latin is None:
        latin = etree.SubElement(rpr, qn("a:latin"))
    latin.set("typeface", typeface)


def add_textbox(
    slide: Slide,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    *,
    anchor: str = "top",
) -> Shape:
    """Add a textbox at inch coordinates and return the shape.

    `anchor` is one of "top", "middle", "bottom".
    """
    tb = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    # python-pptx defaults add_textbox to spAutoFit, which DESIGN.md section
    # 12.1 rule 11 forbids (a textbox autosizing rebreaks layout). Force NONE.
    tf.auto_size = MSO_AUTO_SIZE.NONE
    if anchor == "top":
        tf.vertical_anchor = MSO_ANCHOR.TOP
    elif anchor == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    return tb


def fill_slide_background(slide: Slide, color_hex: str) -> None:
    """Solid background color via the slide's background fill."""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _hex_to_rgb(color_hex)


def fill_slide_background_gradient(slide: Slide, grad: Gradient) -> None:
    """Gradient background via a full-slide rectangle (most reliable across
    PowerPoint and Slides). The bg.fill API does not support gradient fills
    directly, so we draw a rectangle that matches the slide.
    """
    fill_slide_background(slide, INK)  # safety floor
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Emu(int(SLIDE_WIDTH_IN * 914400)), Emu(int(SLIDE_HEIGHT_IN * 914400))
    )
    rect.line.fill.background()
    sppr = rect._element.spPr
    _replace_fill(sppr, _make_grad_fill_xml(grad))


def add_brand_stamp(slide: Slide, on_dark_surface: bool) -> Shape:
    """Top-right wordmark per reference.md global pattern. Always 0.5 in tall.

    on_dark_surface=True selects the white wordmark; False selects gradient.
    """
    asset_path = assets.WHITE_LOGO_PNG() if on_dark_surface else assets.GRADIENT_LOGO_PNG()
    height = Inches(LOGO_STAMP_HEIGHT_IN)
    pic = slide.shapes.add_picture(str(asset_path), Inches(0), Inches(0), height=height)
    width_in = pic.width / 914400
    pic.left = Inches(SLIDE_WIDTH_IN - LOGO_STAMP_RIGHT_IN - width_in)
    pic.top = Inches(LOGO_STAMP_TOP_IN)
    pic.name = "Brand stamp (wordmark)"
    try:
        pic._element._nvXxPr.cNvPr.set("descr", "Reeinvent wordmark")  # type: ignore[attr-defined]
    except Exception:
        pass
    return pic


def add_centered_logo(
    slide: Slide,
    *,
    on_dark_surface: bool,
    height_in: float,
    top_in: float,
) -> Shape:
    """Centered wordmark used on cover and closing slides."""
    asset_path = assets.WHITE_LOGO_PNG() if on_dark_surface else assets.GRADIENT_LOGO_PNG()
    height = Inches(height_in)
    pic = slide.shapes.add_picture(str(asset_path), Inches(0), Inches(top_in), height=height)
    width_in = pic.width / 914400
    pic.left = Inches((SLIDE_WIDTH_IN - width_in) / 2)
    pic.name = "Centered wordmark"
    return pic


def add_arrow_watermark(
    slide: Slide,
    *,
    on_dark_surface: bool,
    size_in: float = 4.5,
) -> Shape:
    """Arrow watermark flush at top-right. Per DESIGN.md section 4 position rule.

    Always fully visible (no bleed off-canvas). One per slide.
    """
    pic = slide.shapes.add_picture(
        str(assets.ARROW_UP_PNG()),
        Inches(0),
        Inches(0),
        width=Inches(size_in),
        height=Inches(size_in),
    )
    pic.left = Inches(SLIDE_WIDTH_IN - size_in)
    pic.top = Inches(0)
    pic.name = "Arrow watermark"
    # Opacity per DESIGN.md section 4: 6-10% on dark, 4-6% on light. Pick the
    # high end so the signature mark reads on small renders / projector dim
    # rooms without overpowering the foreground content.
    alpha_pct = 10 if on_dark_surface else 6
    blip_fill = pic._element.find(qn("p:blipFill"))
    if blip_fill is not None:
        blip = blip_fill.find(qn("a:blip"))
        if blip is not None:
            for existing in blip.findall(qn("a:alphaModFix")):
                blip.remove(existing)
            etree.SubElement(blip, qn("a:alphaModFix"), {"amt": str(alpha_pct * 1000)})
    try:
        pic._element._nvXxPr.cNvPr.set("descr", "Reeinvent arrow watermark")  # type: ignore[attr-defined]
    except Exception:
        pass
    return pic


def add_gradient_stripe(
    slide: Slide,
    *,
    left_in: float,
    top_in: float,
    width_in: float,
    height_pt: float = 2.5,
    grad: Gradient = SIGNATURE_GRADIENT,
) -> Shape:
    """Thin gradient bar (eyebrow underline). Always tracks the text it sits under."""
    h = Pt(height_pt)
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left_in), Inches(top_in), Inches(width_in), h
    )
    rect.line.fill.background()
    fill_shape_with_gradient(rect, grad)
    rect.name = "Gradient stripe"
    return rect


def add_gradient_underline(
    slide: Slide,
    *,
    left_in: float,
    top_in: float,
    width_in: float,
    grad: Gradient = SIGNATURE_GRADIENT,
) -> Shape:
    """Single per-slide gradient underline under emphasized headline word(s).
    DESIGN.md section 8: 2.5 pt thick, sits just below text baseline.
    """
    return add_gradient_stripe(
        slide,
        left_in=left_in,
        top_in=top_in,
        width_in=width_in,
        height_pt=2.5,
        grad=grad,
    )


def add_gradient_pill(
    slide: Slide,
    *,
    text: str,
    left_in: float,
    top_in: float,
    height_in: float = 0.33,
    padding_in: float = 0.16,
    grad: Gradient = SIGNATURE_GRADIENT,
    text_color_hex: str = WHITE,
) -> Shape:
    """Gradient-filled rounded pill. Used for CHALLENGE/SOLUTIONS/RESULTS,
    eyebrows-on-cards, WHAT'S INCLUDED labels.
    """
    role = TYPE_SCALE["eyebrow"]
    estimated_text_width_in = max(0.6, len(text) * role.size_pt * 0.0085)
    width_in = estimated_text_width_in + padding_in * 2
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left_in),
        Inches(top_in),
        Inches(width_in),
        Inches(height_in),
    )
    pill.line.fill.background()
    pill.adjustments[0] = 0.5  # full-pill radius
    fill_shape_with_gradient(pill, grad)
    tf = pill.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run_typography(run, text.upper(), role, color_hex=text_color_hex, weight_override=700)
    pill.name = f"Gradient pill: {text[:20]}"
    return pill


def add_outlined_pill(
    slide: Slide,
    *,
    text: str,
    left_in: float,
    top_in: float,
    height_in: float = 0.33,
    stroke_color_hex: str = WHITE,
    text_color_hex: str = WHITE,
) -> Shape:
    """Outlined (transparent fill) rounded pill. Used for duration pills on
    service slides. Sits on gradient backgrounds, stroke is white.
    """
    role = TYPE_SCALE["eyebrow"]
    estimated_text_width_in = max(0.6, len(text) * role.size_pt * 0.0085)
    width_in = estimated_text_width_in + 0.32
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left_in),
        Inches(top_in),
        Inches(width_in),
        Inches(height_in),
    )
    pill.adjustments[0] = 0.5
    pill.fill.background()  # transparent
    pill.line.color.rgb = _hex_to_rgb(stroke_color_hex)
    pill.line.width = Pt(1.5)
    tf = pill.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run_typography(run, text, role, color_hex=text_color_hex, weight_override=500)
    pill.name = f"Outlined pill: {text[:20]}"
    return pill


def add_card(
    slide: Slide,
    *,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    fill_color_hex: str = WHITE,
    radius_pt: int = 16,
    shadow: bool = True,
) -> Shape:
    """White rounded card with brand-correct Ink-tinted shadow.

    DESIGN.md card spec: 0pt x 4pt, 16pt blur, Ink at 6-10% alpha. Never
    black; always Ink-tinted. We write the shadow as <a:effectLst><a:outerShdw/>
    inside the shape's spPr because python-pptx's high-level shadow API does
    not support setting a custom color or opacity reliably across versions.
    """
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left_in),
        Inches(top_in),
        Inches(width_in),
        Inches(height_in),
    )
    min_dim_in = min(width_in, height_in)
    radius_in = radius_pt / 72
    adj = max(0.02, min(0.5, radius_in / min_dim_in))
    card.adjustments[0] = adj
    card.fill.solid()
    card.fill.fore_color.rgb = _hex_to_rgb(fill_color_hex)
    card.line.fill.background()

    if shadow:
        _apply_brand_shadow(card)

    return card


def _apply_brand_shadow(shape: Shape) -> None:
    """Ink-tinted outer drop shadow per DESIGN.md card spec.

    Implementation: write <a:effectLst><a:outerShdw blurRad="..." dist="..."
    dir="5400000" algn="ctr" rotWithShape="0"><a:srgbClr val="0A1220"><a:alpha
    val="8000"/></a:srgbClr></a:outerShdw></a:effectLst> directly into spPr.
    Direction 5400000 (in 60000ths of a degree) = 90 deg = straight down.
    Distance 50800 EMU = 4 pt. Blur 203200 EMU = 16 pt.
    """
    sppr = shape._element.spPr
    # Remove any existing effectLst (idempotent).
    for child in list(sppr):
        if child.tag == qn("a:effectLst"):
            sppr.remove(child)
    effect_lst = etree.SubElement(sppr, qn("a:effectLst"))
    outer_shdw = etree.SubElement(
        effect_lst,
        qn("a:outerShdw"),
        {
            "blurRad": "203200",   # 16 pt
            "dist": "50800",        # 4 pt
            "dir": "5400000",       # 90 deg (straight down)
            "algn": "ctr",
            "rotWithShape": "0",
        },
    )
    clr = etree.SubElement(outer_shdw, qn("a:srgbClr"), {"val": INK})
    etree.SubElement(clr, qn("a:alpha"), {"val": "8000"})  # 8% alpha


def add_image_or_placeholder(
    slide: Slide,
    *,
    image_path: str | None,
    placeholder_label: str,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    on_dark_surface: bool = True,
    radius_pt: int = 10,
) -> Shape:
    """Drop an image at the named slot; otherwise draw a labeled placeholder.

    The placeholder uses a thin outlined rounded rectangle so the slot is
    visually obvious to the presenter who needs to swap in real photography
    before sending the deck. Per CLAUDE.md no-half-finished-implementations
    rule, the slot is unmistakably a placeholder, not a partial design.
    """
    from pathlib import Path

    if image_path:
        path = Path(image_path)
        if not path.is_file():
            raise FileNotFoundError(f"image not found: {image_path}")
        pic = slide.shapes.add_picture(
            str(path),
            Inches(left_in),
            Inches(top_in),
            width=Inches(width_in),
            height=Inches(height_in),
        )
        pic.name = f"Image: {path.name}"
        try:
            pic._element._nvXxPr.cNvPr.set("descr", placeholder_label)  # type: ignore[attr-defined]
        except Exception:
            pass
        return pic

    # Placeholder: outlined rounded rect with the label centered.
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left_in),
        Inches(top_in),
        Inches(width_in),
        Inches(height_in),
    )
    min_dim = min(width_in, height_in)
    radius_in = radius_pt / 72
    rect.adjustments[0] = max(0.02, min(0.5, radius_in / min_dim))
    if on_dark_surface:
        rect.fill.solid()
        rect.fill.fore_color.rgb = _hex_to_rgb(INK)
        rect.line.color.rgb = _hex_to_rgb(WHITE)
    else:
        rect.fill.solid()
        rect.fill.fore_color.rgb = _hex_to_rgb(OFF_WHITE)
        rect.line.color.rgb = _hex_to_rgb(INK)
    rect.line.width = Pt(0.75)
    rect.name = f"Placeholder: {placeholder_label}"

    label_color = WHITE if on_dark_surface else INK
    tb = add_textbox(
        slide,
        left_in=left_in,
        top_in=top_in + height_in / 2 - 0.3,
        width_in=width_in,
        height_in=0.6,
        anchor="middle",
    )
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run_typography(
        run, placeholder_label, TYPE_SCALE["secondary_body"], color_hex=label_color
    )
    return rect


def add_chat_bubble(
    slide: Slide,
    *,
    text: str,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float = 0.7,
    text_color_hex: str = INK,
    gradient_text: bool = False,
    tail_position: str = "top",
    tail_offset_in: float = 0.6,
) -> Shape:
    """Rounded chat-bubble card with a triangular tail. DESIGN.md /
    reference.md component spec: white rounded card (radius ~10-12pt) plus a
    small triangle (~14pt) extending from the named edge. Used for taglines
    above service names, client descriptors above success-story content.

    `tail_position` is one of: "top", "bottom", "left", "right".
    `tail_offset_in` is the position along the chosen edge from the
    near corner (left or top).
    """
    # Body card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left_in),
        Inches(top_in),
        Inches(width_in),
        Inches(height_in),
    )
    radius_in = 10 / 72
    min_dim = min(width_in, height_in)
    adj = max(0.05, min(0.5, radius_in / min_dim))
    card.adjustments[0] = adj
    card.fill.solid()
    card.fill.fore_color.rgb = _hex_to_rgb(WHITE)
    card.line.fill.background()
    _apply_brand_shadow(card)
    card.name = "Chat bubble"

    # Triangular tail. Drawn as a built-in ISOCELES_TRIANGLE primitive (points
    # up by default) so every PowerPoint / Keynote / Slides / LibreOffice
    # renderer handles it consistently. Rotation positions the tip toward the
    # named edge.
    tail_size = 0.18  # in (~13pt)
    if tail_position == "top":
        tail_left = left_in + tail_offset_in - tail_size / 2
        tail_top = top_in - tail_size + 0.01  # slight overlap so seam vanishes
        rotation = 0
    elif tail_position == "bottom":
        tail_left = left_in + tail_offset_in - tail_size / 2
        tail_top = top_in + height_in - 0.01
        rotation = 180
    elif tail_position == "left":
        tail_left = left_in - tail_size + 0.01
        tail_top = top_in + tail_offset_in - tail_size / 2
        rotation = 270
    elif tail_position == "right":
        tail_left = left_in + width_in - 0.01
        tail_top = top_in + tail_offset_in - tail_size / 2
        rotation = 90
    else:
        raise ValueError(f"unknown tail_position: {tail_position}")

    tail_shape = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(tail_left),
        Inches(tail_top),
        Inches(tail_size),
        Inches(tail_size),
    )
    tail_shape.rotation = rotation
    tail_shape.fill.solid()
    tail_shape.fill.fore_color.rgb = _hex_to_rgb(WHITE)
    tail_shape.line.fill.background()
    tail_shape.name = "Chat bubble tail"

    # Text inside the bubble (separate textbox so we control margins).
    pad = 0.18
    tb = add_textbox(
        slide,
        left_in=left_in + pad,
        top_in=top_in + 0.05,
        width_in=width_in - pad * 2,
        height_in=height_in - 0.1,
        anchor="middle",
    )
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    set_run_typography(
        run,
        text,
        TYPE_SCALE["secondary_body"],
        color_hex=text_color_hex,
        weight_override=700 if gradient_text else 500,
    )
    if gradient_text:
        fill_run_with_gradient(run, SIGNATURE_GRADIENT)

    return card


def add_arrow_bullet_list(
    slide: Slide,
    *,
    items: Iterable[str],
    left_in: float,
    top_in: float,
    width_in: float,
    item_height_in: float = 0.42,
    text_color_hex: str = INK,
) -> None:
    """Bullet list with the canonical Upwards-Arrow PNG marker.

    DESIGN.md component bullet rules: one marker only, one line per item,
    max 6 items.
    """
    items_list = list(items)
    if len(items_list) > 6:
        raise ValueError("bullet list exceeds the 6-item brand maximum (DESIGN.md)")
    role = TYPE_SCALE["body"]
    marker_size_in = role.size_pt * 1.6 / 72  # 1.6 x cap-height approx
    gap_in = 0.18
    for idx, item in enumerate(items_list):
        row_top = top_in + idx * item_height_in
        # marker
        marker_top = row_top + (item_height_in - marker_size_in) / 2
        slide.shapes.add_picture(
            str(assets.UPWARDS_ARROW_PNG()),
            Inches(left_in),
            Inches(marker_top),
            width=Inches(marker_size_in),
            height=Inches(marker_size_in),
        )
        # label
        text_left = left_in + marker_size_in + gap_in
        text_width = width_in - (marker_size_in + gap_in)
        tb = add_textbox(slide, text_left, row_top, text_width, item_height_in, anchor="middle")
        tb.text_frame.word_wrap = False
        tb.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        set_run_typography(run, item, role, color_hex=text_color_hex)


def set_card_text_normautofit(shape: Shape) -> None:
    """Apply <a:normAutofit/> to a shape's text frame so text shrinks within
    the card's fixed geometry. DESIGN.md section 12.1 rule 11.
    """
    tf = shape.text_frame
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        for child in list(bodyPr):
            if child.tag in (qn("a:spAutoFit"), qn("a:normAutofit")):
                bodyPr.remove(child)
        etree.SubElement(bodyPr, qn("a:normAutofit"))


def recolor_picture(pic: Shape, target_hex: str) -> None:
    """Recolor a brand-mark PNG to a target hex via OOXML duotone effect.

    The brand PNGs ship with a single visible fill (off-white for the arrow,
    white for the white wordmark, gradient for the gradient wordmark). To
    place an arrow on a light background, DESIGN.md section 4 requires the
    arrow fill to switch to Ink. We do this with a duotone where both stops
    are the target color, which effectively replaces all visible pixels.
    """
    blip_fill = pic._element.find(qn("p:blipFill"))
    if blip_fill is None:
        return
    blip = blip_fill.find(qn("a:blip"))
    if blip is None:
        return
    # Remove any existing duotone first (idempotent).
    for d in blip.findall(qn("a:duotone")):
        blip.remove(d)
    duo = etree.SubElement(blip, qn("a:duotone"))
    for stop_alpha in (0, 100):
        clr = etree.SubElement(duo, qn("a:srgbClr"), {"val": target_hex})
        if stop_alpha == 0:
            etree.SubElement(clr, qn("a:lumMod"), {"val": "50000"})
            etree.SubElement(clr, qn("a:lumOff"), {"val": "0"})


def assert_no_em_dash(text: str) -> None:
    """DESIGN.md rule 31: zero U+2014 anywhere in any output."""
    if "\u2014" in text:
        raise ValueError(
            f"em-dash (U+2014) found in text: {text!r}. "
            "Replace with hyphen, colon, semicolon, or period (DESIGN.md rule 31)."
        )
