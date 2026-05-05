"""Success story slide. reference.md A6.

Top white strip = client name + SUCCESS STORY label + brand stamp.
Main area = Ink fill with three rows: CHALLENGE / SOLUTION / RESULTS, each
row a gradient pill label + white card. RESULTS card holds an arrow-bullet
list of metrics. Right-side mockup is a placeholder rectangle.
"""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from reeinvent_pitch_deck import assets
from reeinvent_pitch_deck.helpers import (
    add_arrow_bullet_list,
    add_card,
    add_chat_bubble,
    add_gradient_pill,
    add_image_or_placeholder,
    add_textbox,
    fill_run_with_gradient,
    fill_slide_background,
    set_run_typography,
)
from reeinvent_pitch_deck.spec import SlideSpec
from reeinvent_pitch_deck.theme import (
    CORE_BLUE,
    INK,
    PALETTE,
    SIGNATURE_GRADIENT,
    SLIDE_HEIGHT_IN,
    SLIDE_WIDTH_IN,
    TYPE_SCALE,
    TypeRole,
    WHITE,
)


def build(slide: Slide, spec: SlideSpec) -> None:
    fill_slide_background(slide, INK)

    # Top white strip.
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_WIDTH_IN), Inches(0.9)
    )
    strip.line.fill.background()
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor.from_string(WHITE)

    # Client badge (left).
    client_color = PALETTE.get(spec.get("client_color", "core_blue"), CORE_BLUE)
    client_name = spec.get("client_name")
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(0.18), Inches(2.6), Inches(0.55)
    )
    badge.adjustments[0] = 0.5
    badge.line.fill.background()
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor.from_string(client_color)
    btf = badge.text_frame
    btf.margin_left = Pt(0)
    btf.margin_right = Pt(0)
    btf.margin_top = Pt(0)
    btf.margin_bottom = Pt(0)
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    set_run_typography(
        br, client_name.upper(), TYPE_SCALE["body"], color_hex=WHITE, weight_override=700
    )

    # Center label.
    lb = add_textbox(slide, 3.5, 0.18, 6.0, 0.55, anchor="middle")
    lp = lb.text_frame.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    lr = lp.add_run()
    set_run_typography(
        lr,
        "SUCCESS STORY",
        TypeRole(28, 300, 1.05),
        color_hex=INK,
        weight_override=300,
    )
    rpr = lr._r.get_or_add_rPr()
    rpr.set("spc", "200")

    # Right wordmark.
    pic = slide.shapes.add_picture(
        str(assets.GRADIENT_LOGO_PNG()),
        Inches(0),
        Inches(0.2),
        height=Inches(0.5),
    )
    pic.left = Inches(SLIDE_WIDTH_IN - 0.4 - (pic.width / 914400))
    pic.top = Inches(0.2)

    # Three-row content stack at left, ~50% width.
    rows_left = 0.4
    rows_width = 6.5
    rows_top = 1.4
    row_height = 1.5

    # Optional chat bubble at top with triangular tail pointing UP at the
    # client logo on the white strip above. reference.md A6 component spec.
    chat = spec.get("chat_bubble")
    if chat:
        add_chat_bubble(
            slide,
            text=chat,
            left_in=rows_left,
            top_in=rows_top,
            width_in=rows_width,
            height_in=0.7,
            text_color_hex=INK,
            gradient_text=True,
            tail_position="top",
            tail_offset_in=0.7,
        )
        rows_top += 0.95

    rows = [
        ("CHALLENGE", spec.get("challenge")),
        ("SOLUTION", spec.get("solution")),
        ("RESULTS", spec.get("results")),
    ]

    for i, (label, content) in enumerate(rows):
        rt = rows_top + i * (row_height + 0.15)
        # Pill on left.
        add_gradient_pill(
            slide,
            text=label,
            left_in=rows_left,
            top_in=rt + 0.08,
            height_in=0.36,
        )
        # White card on right of pill.
        card_left = rows_left + 1.6
        card_width = rows_width - 1.6
        add_card(
            slide,
            left_in=card_left,
            top_in=rt,
            width_in=card_width,
            height_in=row_height - 0.05,
            fill_color_hex=WHITE,
            radius_pt=12,
        )
        if isinstance(content, list):
            # RESULTS: bullet list of metrics.
            add_arrow_bullet_list(
                slide,
                items=content,
                left_in=card_left + 0.25,
                top_in=rt + 0.15,
                width_in=card_width - 0.5,
                item_height_in=(row_height - 0.4) / max(1, len(content)),
                text_color_hex=INK,
            )
        else:
            tb = add_textbox(
                slide,
                card_left + 0.25,
                rt + 0.05,
                card_width - 0.5,
                row_height - 0.15,
                anchor="middle",
            )
            tb.text_frame.word_wrap = True
            tp = tb.text_frame.paragraphs[0]
            tp.alignment = PP_ALIGN.LEFT
            tr = tp.add_run()
            set_run_typography(tr, content, TYPE_SCALE["secondary_body"], color_hex=INK)

    # Right-side product image or placeholder.
    add_image_or_placeholder(
        slide,
        image_path=spec.get("image"),
        placeholder_label="[PRODUCT SCREENSHOT]",
        left_in=7.4,
        top_in=1.4,
        width_in=SLIDE_WIDTH_IN - 7.4 - 0.4,
        height_in=SLIDE_HEIGHT_IN - 1.4 - 0.6,
        on_dark_surface=True,
    )
