"""Agenda / Table of Contents. DESIGN.md section 6.9.

Off-White. Left column: AGENDA eyebrow + gradient stripe + title.
Right column: numbered list. Each row = large gradient number + label + description.
"""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Pt

from reeinvent_pitch_deck.helpers import (
    add_brand_stamp,
    add_gradient_stripe,
    add_textbox,
    fill_run_with_gradient,
    fill_slide_background,
    set_run_typography,
)
from reeinvent_pitch_deck.spec import SlideSpec
from reeinvent_pitch_deck.theme import (
    INK,
    OFF_WHITE,
    SIGNATURE_GRADIENT,
    SLIDE_HEIGHT_IN,
    SLIDE_WIDTH_IN,
    TYPE_SCALE,
    TypeRole,
)


def build(slide: Slide, spec: SlideSpec) -> None:
    fill_slide_background(slide, OFF_WHITE)
    add_brand_stamp(slide, on_dark_surface=False)

    eyebrow = spec.get("eyebrow", "AGENDA")
    title = spec.get("title", "Today's agenda")

    # Left column: eyebrow + stripe + title.
    eb = add_textbox(slide, 0.6, 0.6, 4.0, 0.4)
    p = eb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    set_run_typography(run, eyebrow.upper(), TYPE_SCALE["eyebrow"], color_hex=INK)
    fill_run_with_gradient(run, SIGNATURE_GRADIENT)
    rpr = run._r.get_or_add_rPr()
    rpr.set("spc", "150")

    add_gradient_stripe(
        slide,
        left_in=0.6,
        top_in=0.92,
        width_in=max(0.6, len(eyebrow) * 0.11),
    )

    tb = add_textbox(slide, 0.6, 1.5, 5.5, 5.0, anchor="top")
    tb.text_frame.word_wrap = True
    tp = tb.text_frame.paragraphs[0]
    tp.alignment = PP_ALIGN.LEFT
    tr = tp.add_run()
    set_run_typography(tr, title, TYPE_SCALE["title"], color_hex=INK)

    # Right column: numbered list. Layout the list anchored to the bottom 2/3
    # of the canvas (sparse-content rule) and scale row height by item count.
    items = spec.get("items")
    list_left = 7.0
    available_h = SLIDE_HEIGHT_IN - 1.2 - 0.6  # below header band, above footer
    row_height = available_h / max(1, len(items))
    list_top = 1.2

    # Adapt typography to row height so titles + descriptions never clip.
    if row_height >= 1.25:
        title_size = 28
        desc_size = 14
    elif row_height >= 0.95:
        title_size = 22
        desc_size = 13
    else:
        title_size = 18
        desc_size = 12

    title_role = TypeRole(title_size, 700, 1.1)
    desc_role = TypeRole(desc_size, 400, 1.3)
    number_role = TypeRole(36, 700, 1.0)

    for i, item in enumerate(items):
        row_top = list_top + i * row_height
        # Number.
        nb = add_textbox(slide, list_left, row_top, 0.9, row_height, anchor="middle")
        np = nb.text_frame.paragraphs[0]
        np.alignment = PP_ALIGN.LEFT
        nr = np.add_run()
        set_run_typography(nr, f"{i + 1:02d}", number_role, color_hex=INK)
        fill_run_with_gradient(nr, SIGNATURE_GRADIENT)

        # Label + description.
        lb = add_textbox(
            slide,
            list_left + 1.0,
            row_top,
            SLIDE_WIDTH_IN - (list_left + 1.0) - 0.4,
            row_height,
            anchor="middle",
        )
        lb.text_frame.word_wrap = True
        lp = lb.text_frame.paragraphs[0]
        lp.alignment = PP_ALIGN.LEFT
        lr = lp.add_run()
        set_run_typography(
            lr, item["title"], title_role, color_hex=INK, weight_override=700
        )
        if item.get("description"):
            dp = lb.text_frame.add_paragraph()
            dp.alignment = PP_ALIGN.LEFT
            dp.space_before = Pt(2)
            dr = dp.add_run()
            set_run_typography(dr, item["description"], desc_role, color_hex=INK)
