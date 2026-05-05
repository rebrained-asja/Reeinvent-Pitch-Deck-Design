"""Quote slide. DESIGN.md section 6.7.

Deep Navy or Off-White background. Opening 140 pt gradient-text quote glyph
top-left. Quote body 40 pt italic centered. Attribution eyebrow-style below.
"""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Pt

from reeinvent_pitch_deck.helpers import (
    add_brand_stamp,
    add_textbox,
    fill_run_with_gradient,
    fill_slide_background,
    set_run_typography,
)
from reeinvent_pitch_deck.spec import SlideSpec
from reeinvent_pitch_deck.theme import (
    DEEP_NAVY,
    INK,
    OFF_WHITE,
    SIGNATURE_GRADIENT,
    SLIDE_HEIGHT_IN,
    SLIDE_WIDTH_IN,
    TYPE_SCALE,
    WHITE,
)


def build(slide: Slide, spec: SlideSpec) -> None:
    bg = spec.get("background", "deep_navy")
    on_dark = bg in {"deep_navy", "ink"}
    text_color = WHITE if on_dark else INK
    fill_slide_background(slide, DEEP_NAVY if bg == "deep_navy" else (INK if bg == "ink" else OFF_WHITE))
    add_brand_stamp(slide, on_dark_surface=on_dark)

    # Opening quote glyph at 140 pt with gradient text fill.
    glyph_box = add_textbox(slide, 0.6, 0.6, 2.5, 2.0)
    p = glyph_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    from reeinvent_pitch_deck.theme import TypeRole
    glyph_role = TypeRole(140, 700, 1.0)
    set_run_typography(run, '"', glyph_role, color_hex=text_color)
    fill_run_with_gradient(run, SIGNATURE_GRADIENT)

    # Quote body.
    quote_text = spec.get("quote")
    qb = add_textbox(
        slide,
        left_in=1.5,
        top_in=2.6,
        width_in=SLIDE_WIDTH_IN - 3.0,
        height_in=3.0,
        anchor="middle",
    )
    qb.text_frame.word_wrap = True
    qp = qb.text_frame.paragraphs[0]
    qp.alignment = PP_ALIGN.CENTER
    qr = qp.add_run()
    set_run_typography(qr, quote_text, TYPE_SCALE["quote"], color_hex=text_color, italic=True)

    # Attribution.
    name = spec.get("attribution_name")
    role = spec.get("attribution_role")
    ab = add_textbox(slide, 1.5, SLIDE_HEIGHT_IN - 1.6, SLIDE_WIDTH_IN - 3.0, 0.6)
    ap = ab.text_frame.paragraphs[0]
    ap.alignment = PP_ALIGN.CENTER
    nr = ap.add_run()
    set_run_typography(nr, name.upper(), TYPE_SCALE["eyebrow"], color_hex=text_color, weight_override=500)
    fill_run_with_gradient(nr, SIGNATURE_GRADIENT)
    rpr = nr._r.get_or_add_rPr()
    rpr.set("spc", "150")
    if role:
        sep = ap.add_run()
        set_run_typography(sep, "  ", TYPE_SCALE["eyebrow"], color_hex=text_color)
        rr = ap.add_run()
        set_run_typography(rr, role.upper(), TYPE_SCALE["eyebrow"], color_hex=text_color)
        rpr2 = rr._r.get_or_add_rPr()
        rpr2.set("spc", "150")
