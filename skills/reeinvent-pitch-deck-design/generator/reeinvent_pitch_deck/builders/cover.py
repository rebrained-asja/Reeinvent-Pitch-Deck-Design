"""Cover slide. DESIGN.md section 6.1 + reference.md A1.

Two cover variants the spec can pick between via `background`:
  - "ink" (DESIGN.md 6.1): Ink solid bg, eyebrow top-left, title left-aligned,
    sub-headline, white wordmark bottom-left, date bottom-right.
  - "gradient" (reference.md A1): Signature Gradient full-bleed, large arrow
    watermark top-right at 8% opacity, REEINVENT wordmark centered.
"""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from reeinvent_pitch_deck.helpers import (
    add_arrow_watermark,
    add_brand_stamp,
    add_centered_logo,
    add_gradient_underline,
    add_textbox,
    fill_run_with_gradient,
    fill_slide_background,
    fill_slide_background_gradient,
    set_run_typography,
)
from reeinvent_pitch_deck.spec import SlideSpec
from reeinvent_pitch_deck.theme import (
    INK,
    OFF_WHITE,
    SIGNATURE_GRADIENT,
    SLIDE_HEIGHT_IN,
    SLIDE_WIDTH_IN,
    TYPE_SCALE,
    WHITE,
)


def build(slide: Slide, spec: SlideSpec) -> None:
    bg = spec.get("background", "ink")
    title = spec.get("title")
    subtitle = spec.get("subtitle")
    eyebrow = spec.get("eyebrow")
    highlight_words = spec.get("highlight_words", [])
    date = spec.get("date")
    author = spec.get("author")

    if bg == "gradient":
        _build_gradient_cover(slide, title, subtitle)
    else:
        _build_ink_cover(slide, title, subtitle, eyebrow, highlight_words, date, author)


def _build_gradient_cover(slide: Slide, title: str, subtitle: str | None) -> None:
    fill_slide_background_gradient(slide, SIGNATURE_GRADIENT)
    add_arrow_watermark(slide, on_dark_surface=True, size_in=4.5)
    # White wordmark centered horizontally, anchored vertical center minus 0.3 in.
    # Logo natural width is large; centered_logo handles math.
    logo_height = 0.9
    add_centered_logo(
        slide,
        on_dark_surface=True,
        height_in=logo_height,
        top_in=(SLIDE_HEIGHT_IN - logo_height) / 2 - 0.3,
    )
    # Optional subtitle line below the centered wordmark.
    if title or subtitle:
        tb = add_textbox(
            slide,
            left_in=0.6,
            top_in=SLIDE_HEIGHT_IN - 1.4,
            width_in=SLIDE_WIDTH_IN - 1.2,
            height_in=1.0,
            anchor="bottom",
        )
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        if title:
            r = p.add_run()
            set_run_typography(r, title, TYPE_SCALE["sub_headline"], color_hex=WHITE)
        if subtitle:
            tb.text_frame.add_paragraph().alignment = PP_ALIGN.CENTER
            p2 = tb.text_frame.paragraphs[-1]
            r2 = p2.add_run()
            set_run_typography(r2, subtitle, TYPE_SCALE["body"], color_hex=WHITE)


def _build_ink_cover(
    slide: Slide,
    title: str,
    subtitle: str | None,
    eyebrow: str | None,
    highlight_words: list[str],
    date: str | None,
    author: str | None,
) -> None:
    fill_slide_background(slide, INK)
    add_arrow_watermark(slide, on_dark_surface=True, size_in=4.5)

    # Eyebrow at (0.8, 0.8). DESIGN.md rule 23: under 40 pt on dark surfaces
    # is WHITE (gradient text forbidden). Gradient stripe sits below the
    # eyebrow per DESIGN.md section 8 stripe rule (width tracks text).
    if eyebrow:
        eb = add_textbox(slide, 0.8, 0.8, SLIDE_WIDTH_IN - 1.6, 0.4)
        p = eb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        set_run_typography(
            run, eyebrow.upper(), TYPE_SCALE["eyebrow"], color_hex=WHITE
        )
        rpr = run._r.get_or_add_rPr()
        rpr.set("spc", "150")
        from reeinvent_pitch_deck.helpers import add_gradient_stripe
        approx_width_in = max(0.6, len(eyebrow) * TYPE_SCALE["eyebrow"].size_pt * 0.011)
        add_gradient_stripe(
            slide, left_in=0.8, top_in=1.12, width_in=approx_width_in
        )

    # Title at left margin, anchored bottom-leftish.
    # DESIGN.md section 6.1 + section 3 + section 8: cover title carries ONE
    # gradient highlight (rectangle behind one or two key words, white text on
    # top). This is the dramatic Reeinvent cover treatment.
    title_top = (SLIDE_HEIGHT_IN / 2) - 1.5
    role = TYPE_SCALE["mega_title"]

    # Estimate title text width & line break to position the gradient highlight
    # rectangle. We approximate per-character width at the mega-title size.
    char_w_in = role.size_pt * 0.0085

    # Render title with gradient text fill on the highlighted phrase. The
    # mega-title size (80 pt) is well above the 40-pt threshold in rule 23, so
    # gradient text is permitted. Other words stay white.
    tb = add_textbox(slide, 0.8, title_top, 11.7, 3.0, anchor="top")
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    if highlight_words:
        remaining = title
        for hl in highlight_words:
            i = remaining.find(hl)
            if i == -1:
                continue
            before = remaining[:i]
            if before:
                r = p.add_run()
                set_run_typography(r, before, role, color_hex=WHITE)
            r = p.add_run()
            set_run_typography(r, hl, role, color_hex=WHITE)
            fill_run_with_gradient(r, SIGNATURE_GRADIENT)
            remaining = remaining[i + len(hl):]
        if remaining:
            r = p.add_run()
            set_run_typography(r, remaining, role, color_hex=WHITE)
    else:
        run = p.add_run()
        set_run_typography(run, title, role, color_hex=WHITE)

    if subtitle:
        sp = tb.text_frame.add_paragraph()
        sp.alignment = PP_ALIGN.LEFT
        sp.space_before = Pt(24)
        r = sp.add_run()
        set_run_typography(r, subtitle, TYPE_SCALE["sub_headline"], color_hex=OFF_WHITE)

    # Gradient underline (one per slide, DESIGN.md rule 9) on a NON-highlighted
    # word - rule 11 forbids stacking gradient text + underline on the same
    # word. Default: place under the title's first word.
    first_word = title.split(" ")[0] if title else ""
    if first_word:
        underline_w_in = min(len(first_word) * char_w_in, 4.5)
        underline_top = title_top + (role.size_pt / 72) + 0.05
        add_gradient_underline(
            slide,
            left_in=0.8,
            top_in=underline_top,
            width_in=underline_w_in,
        )

    # White wordmark bottom-left.
    from reeinvent_pitch_deck import assets
    pic = slide.shapes.add_picture(
        str(assets.WHITE_LOGO_PNG()),
        Inches(0.5),
        Inches(0),
        height=Inches(0.8),
    )
    pic.top = Inches(SLIDE_HEIGHT_IN - 0.5 - (pic.height / 914400))

    # Date / author bottom-right.
    if date or author:
        tag = " ".join(filter(None, [date, author])).strip()
        info = add_textbox(slide, SLIDE_WIDTH_IN - 4.5, SLIDE_HEIGHT_IN - 0.7, 4.0, 0.4)
        ip = info.text_frame.paragraphs[0]
        ip.alignment = PP_ALIGN.RIGHT
        ir = ip.add_run()
        set_run_typography(ir, tag, TYPE_SCALE["footer"], color_hex=OFF_WHITE)
