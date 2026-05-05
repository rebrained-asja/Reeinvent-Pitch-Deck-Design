"""Single-column content slide. DESIGN.md section 6.3.

Off-White background. Eyebrow + gradient stripe at top-left. Title 54 pt Ink
left-aligned. Body 20 pt Ink OR an arrow-bullet list (DESIGN.md bullet rules).
Footer: gradient wordmark + page number.
"""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Pt

from reeinvent_pitch_deck.helpers import (
    add_arrow_bullet_list,
    add_brand_stamp,
    add_gradient_stripe,
    add_gradient_underline,
    add_textbox,
    fill_run_with_gradient,
    fill_slide_background,
    fill_slide_background_gradient,
    set_run_typography,
)
from reeinvent_pitch_deck.spec import SlideSpec
from reeinvent_pitch_deck.theme import (
    DARK_GRADIENT,
    DEEP_NAVY,
    INK,
    OFF_WHITE,
    SIGNATURE_GRADIENT,
    SLIDE_HEIGHT_IN,
    SLIDE_WIDTH_IN,
    TYPE_SCALE,
    VIVID_GRADIENT,
    WHITE,
)


def build(slide: Slide, spec: SlideSpec) -> None:
    bg = spec.get("background", "off_white")
    on_dark = bg in {"ink", "deep_navy", "gradient", "dark_gradient", "vivid_gradient"}

    if bg == "off_white":
        fill_slide_background(slide, OFF_WHITE)
    elif bg == "ink":
        fill_slide_background(slide, INK)
    elif bg == "deep_navy":
        fill_slide_background(slide, DEEP_NAVY)
    elif bg in ("gradient", "dark_gradient", "vivid_gradient"):
        grad = {
            "gradient": SIGNATURE_GRADIENT,
            "dark_gradient": DARK_GRADIENT,
            "vivid_gradient": VIVID_GRADIENT,
        }[bg]
        fill_slide_background_gradient(slide, grad)

    add_brand_stamp(slide, on_dark_surface=on_dark)

    text_color = WHITE if on_dark else INK

    eyebrow = spec.get("eyebrow")
    eyebrow_top = 0.6
    title_top = 1.2
    if eyebrow:
        eb = add_textbox(slide, 0.6, eyebrow_top, 8.0, 0.4)
        p = eb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        set_run_typography(
            run, eyebrow.upper(), TYPE_SCALE["eyebrow"], color_hex=text_color
        )
        if not on_dark:
            fill_run_with_gradient(run, SIGNATURE_GRADIENT)
        rpr = run._r.get_or_add_rPr()
        rpr.set("spc", "150")
        # Gradient stripe: width matches text width (DESIGN.md section 8 stripe rule).
        # Approximate from char count.
        approx_width_in = max(0.6, len(eyebrow) * TYPE_SCALE["eyebrow"].size_pt * 0.011)
        add_gradient_stripe(
            slide,
            left_in=0.6,
            top_in=eyebrow_top + 0.32,
            width_in=approx_width_in,
            height_pt=2.5,
        )
        title_top = 1.1

    title = spec.get("title")
    highlight_words = spec.get("highlight_words", [])
    title_box = add_textbox(slide, 0.6, title_top, 11.0, 1.6)
    title_box.text_frame.word_wrap = True
    p = title_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _emit_title_runs(p, title, highlight_words, text_color, TYPE_SCALE["title"])

    if highlight_words:
        first = highlight_words[0]
        approx_w = len(first) * TYPE_SCALE["title"].size_pt * 0.0085 + 0.15
        # Stripe placement is rough; sit ~6 pt below baseline of title row.
        underline_top = title_top + (TYPE_SCALE["title"].size_pt / 72) + 0.05
        add_gradient_underline(
            slide,
            left_in=0.6,
            top_in=underline_top,
            width_in=min(approx_w, 9.0),
        )

    body = spec.get("body")
    bullets = spec.get("bullets")

    body_top = title_top + 1.6 + 0.3
    if bullets:
        add_arrow_bullet_list(
            slide,
            items=bullets,
            left_in=0.6,
            top_in=body_top,
            width_in=11.0,
            text_color_hex=text_color,
        )
    elif body:
        body_box = add_textbox(slide, 0.6, body_top, 9.5, 4.0)
        body_box.text_frame.word_wrap = True
        bp = body_box.text_frame.paragraphs[0]
        bp.alignment = PP_ALIGN.LEFT
        run = bp.add_run()
        set_run_typography(run, body, TYPE_SCALE["body"], color_hex=text_color)


def _emit_title_runs(paragraph, title: str, highlights: list[str], color: str, role) -> None:
    if not highlights:
        run = paragraph.add_run()
        set_run_typography(run, title, role, color_hex=color)
        return
    remaining = title
    for hl in highlights:
        idx = remaining.find(hl)
        if idx == -1:
            continue
        before = remaining[:idx]
        if before:
            run = paragraph.add_run()
            set_run_typography(run, before, role, color_hex=color)
        run = paragraph.add_run()
        set_run_typography(run, hl, role, color_hex=color)
        fill_run_with_gradient(run, SIGNATURE_GRADIENT)
        remaining = remaining[idx + len(hl):]
    if remaining:
        run = paragraph.add_run()
        set_run_typography(run, remaining, role, color_hex=color)
