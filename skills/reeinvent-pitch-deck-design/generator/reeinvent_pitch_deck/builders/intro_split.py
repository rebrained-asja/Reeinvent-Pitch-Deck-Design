"""Intro split slide. reference.md A3.

50/50 vertical split. Left panel = top half image-placeholder + bottom half
gradient block with two-line display title (thin word over heavy word).
Right panel = Ink fill, headline + body, brand stamp.
"""

from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from reeinvent_pitch_deck.helpers import (
    add_arrow_watermark,
    add_brand_stamp,
    add_gradient_underline,
    add_image_or_placeholder,
    add_textbox,
    fill_run_with_gradient,
    fill_shape_with_gradient,
    set_run_typography,
)
from reeinvent_pitch_deck.spec import SlideSpec
from reeinvent_pitch_deck.theme import (
    DEEP_NAVY,
    INK,
    SIGNATURE_GRADIENT,
    SLIDE_HEIGHT_IN,
    SLIDE_WIDTH_IN,
    TypeRole,
    WHITE,
)
from pptx.dml.color import RGBColor


def build(slide: Slide, spec: SlideSpec) -> None:
    # Right panel: Ink fill (full slide first, left panel painted on top).
    from reeinvent_pitch_deck.helpers import fill_slide_background
    fill_slide_background(slide, INK)
    add_arrow_watermark(slide, on_dark_surface=True, size_in=2.0)
    add_brand_stamp(slide, on_dark_surface=True)

    # Left panel: 50% width, top half is placeholder rectangle (Deep Navy),
    # bottom half is gradient block with the split title.
    left_panel_w = SLIDE_WIDTH_IN / 2
    half_h = SLIDE_HEIGHT_IN / 2

    # Top half: team photo OR placeholder.
    image_note = spec.get("image_note", "[TEAM PHOTO]")
    add_image_or_placeholder(
        slide,
        image_path=spec.get("image"),
        placeholder_label=image_note,
        left_in=0,
        top_in=0,
        width_in=left_panel_w,
        height_in=half_h,
        on_dark_surface=True,
        radius_pt=0,
    )

    # Bottom half: signature gradient.
    grad_block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(half_h),
        Inches(left_panel_w),
        Inches(half_h),
    )
    grad_block.line.fill.background()
    fill_shape_with_gradient(grad_block, SIGNATURE_GRADIENT)

    # Two-line split title inside the gradient block.
    tb = add_textbox(slide, 0.6, half_h + 0.4, left_panel_w - 1.2, half_h - 0.8, anchor="top")
    tb.text_frame.word_wrap = True
    # Reference.md A3: top word ~72 pt thin, bottom word ~120 pt black.
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    set_run_typography(
        r1,
        spec.get("split_top"),
        TypeRole(72, 300, 1.0),
        color_hex=WHITE,
        weight_override=300,
    )
    p2 = tb.text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    set_run_typography(
        r2,
        spec.get("split_bottom"),
        TypeRole(120, 900, 1.0),
        color_hex=WHITE,
        weight_override=900,
    )

    # Right panel: headline + body.
    headline = spec.get("headline")
    body = spec.get("body")
    highlight_words = spec.get("highlight_words", [])

    headline_left = left_panel_w + 0.6
    headline_top = 1.4
    headline_w = SLIDE_WIDTH_IN - headline_left - 0.6

    hb = add_textbox(slide, headline_left, headline_top, headline_w, 2.4)
    hb.text_frame.word_wrap = True
    hp = hb.text_frame.paragraphs[0]
    hp.alignment = PP_ALIGN.LEFT
    _emit_headline_runs(hp, headline.upper(), highlight_words, WHITE, TypeRole(40, 700, 1.1))

    if highlight_words:
        first = highlight_words[0].upper()
        approx_w = len(first) * 40 * 0.0085 + 0.15
        add_gradient_underline(
            slide,
            left_in=headline_left,
            top_in=headline_top + (40 / 72) + 0.05,
            width_in=min(approx_w, headline_w - 0.5),
        )

    bb = add_textbox(slide, headline_left, headline_top + 2.6, headline_w, 3.0)
    bb.text_frame.word_wrap = True
    bp = bb.text_frame.paragraphs[0]
    bp.alignment = PP_ALIGN.LEFT
    br = bp.add_run()
    set_run_typography(br, body, TypeRole(20, 400, 1.4), color_hex=WHITE)


def _emit_headline_runs(paragraph, headline: str, highlights, color, role):
    if not highlights:
        run = paragraph.add_run()
        set_run_typography(run, headline, role, color_hex=color)
        return
    remaining = headline
    upper_highlights = [h.upper() for h in highlights]
    for hl in upper_highlights:
        idx = remaining.find(hl)
        if idx == -1:
            continue
        before = remaining[:idx]
        if before:
            run = paragraph.add_run()
            set_run_typography(run, before, role, color_hex=color)
        run = paragraph.add_run()
        set_run_typography(run, hl, role, color_hex=color)
        # Headline highlights on dark surfaces use solid white plus the gradient
        # underline. Don't apply gradient text fill here (under 40 pt).
        remaining = remaining[idx + len(hl):]
    if remaining:
        run = paragraph.add_run()
        set_run_typography(run, remaining, role, color_hex=color)
