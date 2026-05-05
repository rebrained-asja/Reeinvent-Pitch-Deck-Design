"""Closing slide. DESIGN.md section 6.8 + reference.md A2 (Thank You).

Full-bleed Signature Gradient. Display arrow centered. Mega message white,
700, centered under arrow. Optional CTA pill (white fill, Ink text). White
wordmark bottom-center.
"""

from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from reeinvent_pitch_deck import assets
from reeinvent_pitch_deck.helpers import (
    add_centered_logo,
    add_textbox,
    fill_slide_background_gradient,
    set_run_typography,
)
from reeinvent_pitch_deck.spec import SlideSpec
from reeinvent_pitch_deck.theme import (
    INK,
    SIGNATURE_GRADIENT,
    SLIDE_HEIGHT_IN,
    SLIDE_WIDTH_IN,
    TYPE_SCALE,
    TypeRole,
    WHITE,
)


def build(slide: Slide, spec: SlideSpec) -> None:
    fill_slide_background_gradient(slide, SIGNATURE_GRADIENT)

    # Display arrow: 2.6 in (DESIGN.md 6.8 calls for 6 in but that overpowers
    # the message; 2.6 in is the production-deck balance that keeps the arrow
    # iconic without dominating the hero phrase).
    arrow_size = 2.6
    pic = slide.shapes.add_picture(
        str(assets.ARROW_UP_PNG()),
        Inches((SLIDE_WIDTH_IN - arrow_size) / 2),
        Inches(1.0),
        width=Inches(arrow_size),
        height=Inches(arrow_size),
    )
    pic.name = "Closing slide hero arrow"

    # Mega message.
    message = spec.get("message")
    msg_role = TypeRole(72, 700, 1.05)
    mb = add_textbox(
        slide,
        left_in=0.6,
        top_in=4.0,
        width_in=SLIDE_WIDTH_IN - 1.2,
        height_in=1.4,
        anchor="middle",
    )
    mb.text_frame.word_wrap = True
    mp = mb.text_frame.paragraphs[0]
    mp.alignment = PP_ALIGN.CENTER
    mr = mp.add_run()
    set_run_typography(mr, message, msg_role, color_hex=WHITE)

    # Optional CTA: pill, white fill, Ink label.
    cta_label = spec.get("cta_label")
    if cta_label:
        pill_h_in = 58 / 72
        pill_w_in = max(2.0, len(cta_label) * 0.16 + 0.8)
        pill_left = (SLIDE_WIDTH_IN - pill_w_in) / 2
        pill_top = 5.6
        pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(pill_left),
            Inches(pill_top),
            Inches(pill_w_in),
            Inches(pill_h_in),
        )
        pill.adjustments[0] = 0.5
        pill.fill.solid()
        from pptx.dml.color import RGBColor
        pill.fill.fore_color.rgb = RGBColor.from_string(WHITE)
        pill.line.fill.background()
        tf = pill.text_frame
        tf.margin_left = Pt(0)
        tf.margin_right = Pt(0)
        tf.margin_top = Pt(0)
        tf.margin_bottom = Pt(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = tf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        set_run_typography(cr, cta_label, TYPE_SCALE["cta_label"], color_hex=INK)

    # Centered white wordmark, sitting flush against the bottom safe-area
    # padding so the CTA pill does not overlap. DESIGN.md section 6.8.
    logo_height = 0.7
    add_centered_logo(
        slide,
        on_dark_surface=True,
        height_in=logo_height,
        top_in=SLIDE_HEIGHT_IN - 0.4 - logo_height,
    )
