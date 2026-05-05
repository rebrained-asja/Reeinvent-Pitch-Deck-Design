"""Three-up card row. DESIGN.md section 6.6.

Off-White. Standard eyebrow + title at top. Three white rounded cards anchored
to the bottom 2/3 of the canvas (sparse-content rule, DESIGN.md section 9).
"""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from reeinvent_pitch_deck import assets
from reeinvent_pitch_deck.helpers import (
    add_brand_stamp,
    add_card,
    add_gradient_stripe,
    add_gradient_underline,
    add_textbox,
    fill_run_with_gradient,
    fill_slide_background,
    recolor_picture,
    set_card_text_normautofit,
    set_run_typography,
)
from reeinvent_pitch_deck.spec import SlideSpec
from reeinvent_pitch_deck.theme import (
    CORE_BLUE,
    INK,
    OFF_WHITE,
    SIGNATURE_GRADIENT,
    SLIDE_HEIGHT_IN,
    SLIDE_WIDTH_IN,
    TYPE_SCALE,
    WHITE,
)


def build(slide: Slide, spec: SlideSpec) -> None:
    fill_slide_background(slide, OFF_WHITE)
    add_brand_stamp(slide, on_dark_surface=False)

    eyebrow = spec.get("eyebrow")
    title = spec.get("title")
    highlight_words = spec.get("highlight_words", [])

    eyebrow_top = 0.6
    title_top = 1.2
    if eyebrow:
        eb = add_textbox(slide, 0.6, eyebrow_top, 8.0, 0.4)
        p = eb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        set_run_typography(run, eyebrow.upper(), TYPE_SCALE["eyebrow"], color_hex=INK)
        fill_run_with_gradient(run, SIGNATURE_GRADIENT)
        rpr = run._r.get_or_add_rPr()
        rpr.set("spc", "150")
        approx_width_in = max(0.6, len(eyebrow) * TYPE_SCALE["eyebrow"].size_pt * 0.011)
        add_gradient_stripe(
            slide,
            left_in=0.6,
            top_in=eyebrow_top + 0.32,
            width_in=approx_width_in,
        )
        title_top = 1.1

    title_box = add_textbox(slide, 0.6, title_top, 11.0, 1.4)
    title_box.text_frame.word_wrap = True
    tp = title_box.text_frame.paragraphs[0]
    tp.alignment = PP_ALIGN.LEFT
    _emit_title_runs(tp, title, highlight_words, INK, TYPE_SCALE["title"])
    if highlight_words:
        first = highlight_words[0]
        approx_w = len(first) * TYPE_SCALE["title"].size_pt * 0.0085 + 0.15
        add_gradient_underline(
            slide,
            left_in=0.6,
            top_in=title_top + (TYPE_SCALE["title"].size_pt / 72) + 0.05,
            width_in=min(approx_w, 9.0),
        )

    # Three cards anchored to the lower portion (DESIGN.md sparse-content rule).
    cards = spec.get("cards")
    card_width = 3.85
    gap = 0.25
    total_w = 3 * card_width + 2 * gap
    cards_left = (SLIDE_WIDTH_IN - total_w) / 2
    cards_top = 3.6
    cards_height = 3.3
    for idx, c in enumerate(cards):
        cx = cards_left + idx * (card_width + gap)
        card = add_card(
            slide,
            left_in=cx,
            top_in=cards_top,
            width_in=card_width,
            height_in=cards_height,
            fill_color_hex=WHITE,
            radius_pt=16,
        )
        # Card title.
        tb = add_textbox(
            slide,
            cx + 0.3,
            cards_top + 0.4,
            card_width - 0.6,
            0.6,
            anchor="top",
        )
        tp = tb.text_frame.paragraphs[0]
        tp.alignment = PP_ALIGN.LEFT
        tr = tp.add_run()
        set_run_typography(
            tr, c["title"], TYPE_SCALE["body"], color_hex=INK, weight_override=700
        )
        # Card body.
        bb = add_textbox(
            slide,
            cx + 0.3,
            cards_top + 1.1,
            card_width - 0.6,
            cards_height - 1.5,
            anchor="top",
        )
        bb.text_frame.word_wrap = True
        bp = bb.text_frame.paragraphs[0]
        bp.alignment = PP_ALIGN.LEFT
        br = bp.add_run()
        set_run_typography(
            br, c["body"], TYPE_SCALE["secondary_body"], color_hex=INK
        )
        set_card_text_normautofit(bb)
        # Small arrow accent bottom-right of card. Core Blue per DESIGN.md
        # section 6.6 ("Optional 'Arrow mark': Small arrow (0.5 in), Core
        # Blue, bottom-right corner of card.").
        arrow_size = 0.5
        arrow_pic = slide.shapes.add_picture(
            str(assets.ARROW_UP_PNG()),
            Inches(cx + card_width - arrow_size - 0.25),
            Inches(cards_top + cards_height - arrow_size - 0.25),
            width=Inches(arrow_size),
            height=Inches(arrow_size),
        )
        recolor_picture(arrow_pic, CORE_BLUE)


def _emit_title_runs(paragraph, title: str, highlights: list[str], color: str, role) -> None:
    if not highlights:
        run = paragraph.add_run()
        set_run_typography(run, title, role, color_hex=color)
        return
    remaining = title
    for hl in highlights:
        idx = remaining.find(hl)
        if idx == -1:
            continue
        before = remaining[:idx]
        if before:
            run = paragraph.add_run()
            set_run_typography(run, before, role, color_hex=color)
        run = paragraph.add_run()
        set_run_typography(run, hl, role, color_hex=color)
        fill_run_with_gradient(run, SIGNATURE_GRADIENT)
        remaining = remaining[idx + len(hl):]
    if remaining:
        run = paragraph.add_run()
        set_run_typography(run, remaining, role, color_hex=color)
