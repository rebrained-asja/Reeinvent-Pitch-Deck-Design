"""Stat slide. DESIGN.md section 6.5.

Center-aligned big number at 120 pt with Signature Gradient text fill (per
section 12.1 rule 12 - mandatory). Stat label uppercase below. Optional
supporting line. Optional small Arrow accent below.
"""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from reeinvent_pitch_deck import assets
from reeinvent_pitch_deck.helpers import (
    add_brand_stamp,
    add_textbox,
    fill_run_with_gradient,
    fill_slide_background,
    recolor_picture,
    set_run_typography,
)
from reeinvent_pitch_deck.spec import SlideSpec
from reeinvent_pitch_deck.theme import (
    CORE_BLUE,
    INK,
    OFF_WHITE,
    SIGNATURE_GRADIENT,
    SLIDE_HEIGHT_IN,
    SLIDE_WIDTH_IN,
    TYPE_SCALE,
    WHITE,
)


def build(slide: Slide, spec: SlideSpec) -> None:
    bg = spec.get("background", "off_white")
    on_dark = bg == "ink"
    text_color = WHITE if on_dark else INK

    fill_slide_background(slide, INK if on_dark else OFF_WHITE)
    add_brand_stamp(slide, on_dark_surface=on_dark)

    number = spec.get("number")
    label = spec.get("label")
    supporting = spec.get("supporting")

    # Stat number: 120 pt, center, Signature Gradient text fill (mandatory).
    num_box = add_textbox(
        slide,
        left_in=1.0,
        top_in=SLIDE_HEIGHT_IN / 3 - 1.0,
        width_in=SLIDE_WIDTH_IN - 2.0,
        height_in=2.0,
        anchor="middle",
    )
    p = num_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run_typography(run, number, TYPE_SCALE["stat_number"], color_hex=text_color)
    fill_run_with_gradient(run, SIGNATURE_GRADIENT)

    # Stat label: 16 pt uppercase +1.5 pt tracking, weight 500.
    label_box = add_textbox(
        slide,
        left_in=1.0,
        top_in=SLIDE_HEIGHT_IN / 3 + 1.4,
        width_in=SLIDE_WIDTH_IN - 2.0,
        height_in=0.4,
    )
    lp = label_box.text_frame.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    lr = lp.add_run()
    set_run_typography(lr, label.upper(), TYPE_SCALE["stat_label"], color_hex=text_color)
    rpr = lr._r.get_or_add_rPr()
    rpr.set("spc", "150")

    # Supporting sentence: 20 pt regular, 40 pt below label.
    if supporting:
        sb = add_textbox(
            slide,
            left_in=2.5,
            top_in=SLIDE_HEIGHT_IN / 3 + 2.2,
            width_in=SLIDE_WIDTH_IN - 5.0,
            height_in=1.2,
        )
        sb.text_frame.word_wrap = True
        sp = sb.text_frame.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        sr = sp.add_run()
        set_run_typography(sr, supporting, TYPE_SCALE["body"], color_hex=text_color)

    # Small arrow accent: 0.6 in, Core Blue tinted, below supporting line.
    arrow_size = 0.6
    pic = slide.shapes.add_picture(
        str(assets.ARROW_UP_PNG()),
        Inches((SLIDE_WIDTH_IN - arrow_size) / 2),
        Inches(SLIDE_HEIGHT_IN - 1.5),
        width=Inches(arrow_size),
        height=Inches(arrow_size),
    )
    pic.name = "Stat slide accent arrow"
    # DESIGN.md section 4 fill rules: arrow on Off-White must be Ink; on Ink must
    # stay native off-white (no recolor needed). On Off-White we want Core Blue
    # for the accent role per section 6.5.
    if not on_dark:
        recolor_picture(pic, CORE_BLUE)
