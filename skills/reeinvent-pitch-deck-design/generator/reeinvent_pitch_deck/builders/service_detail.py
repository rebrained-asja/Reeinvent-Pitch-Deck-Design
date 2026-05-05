"""Service detail slide. reference.md A5 (the offer template).

Top white strip = AI-DRIVEN [CATEGORY] header (left) + brand stamp (right).
Left sidebar (~25% width) = Signature Gradient with arrow watermark, duration
pill, two-line service name (thin over heavy), chat-bubble callout, and tagline.
Main area = Ink fill with WHAT'S INCLUDED card holding bullet list. The
laptop+screenshot mockup cluster is left as a bordered placeholder rectangle.
"""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from reeinvent_pitch_deck import assets
from reeinvent_pitch_deck.helpers import (
    add_arrow_bullet_list,
    add_arrow_watermark,
    add_card,
    add_chat_bubble,
    add_image_or_placeholder,
    add_outlined_pill,
    add_textbox,
    fill_run_with_gradient,
    fill_shape_with_gradient,
    fill_slide_background,
    set_run_typography,
)
from reeinvent_pitch_deck.spec import SlideSpec
from reeinvent_pitch_deck.theme import (
    CORE_BLUE,
    INK,
    SIGNATURE_GRADIENT,
    SLIDE_HEIGHT_IN,
    SLIDE_WIDTH_IN,
    TYPE_SCALE,
    TypeRole,
    WHITE,
)


def build(slide: Slide, spec: SlideSpec) -> None:
    fill_slide_background(slide, INK)

    # Top white strip, full-width, 0.9 in tall.
    strip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_WIDTH_IN), Inches(0.9)
    )
    strip.line.fill.background()
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor.from_string(WHITE)

    # Top-strip header: "AI-DRIVEN [CATEGORY]" pattern.
    category = spec.get("category", "AI-DRIVEN SOLUTIONS").upper()
    parts = category.split(" ", 1)
    head_left = 0.6
    hb = add_textbox(slide, head_left, 0.18, 8.0, 0.6, anchor="middle")
    hp = hb.text_frame.paragraphs[0]
    hp.alignment = PP_ALIGN.LEFT
    role_head = TypeRole(36, 300, 1.05)
    if len(parts) == 2:
        first, second = parts
        r1 = hp.add_run()
        set_run_typography(r1, first + " ", role_head, color_hex=INK, weight_override=300)
        r2 = hp.add_run()
        set_run_typography(
            r2,
            second,
            TypeRole(36, 700, 1.05),
            color_hex=INK,
            weight_override=700,
        )
        fill_run_with_gradient(r2, SIGNATURE_GRADIENT)
    else:
        r1 = hp.add_run()
        set_run_typography(r1, category, role_head, color_hex=INK, weight_override=700)
        fill_run_with_gradient(r1, SIGNATURE_GRADIENT)

    # Top-strip wordmark on the right (gradient on white).
    pic = slide.shapes.add_picture(
        str(assets.GRADIENT_LOGO_PNG()),
        Inches(0),
        Inches(0.2),
        height=Inches(0.5),
    )
    pic.left = Inches(SLIDE_WIDTH_IN - 0.4 - (pic.width / 914400))
    pic.top = Inches(0.2)

    # Left sidebar: gradient fill, ~25% width.
    sidebar_w = SLIDE_WIDTH_IN * 0.28
    sidebar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0.9),
        Inches(sidebar_w),
        Inches(SLIDE_HEIGHT_IN - 0.9),
    )
    sidebar.line.fill.background()
    fill_shape_with_gradient(sidebar, SIGNATURE_GRADIENT)

    # Sidebar arrow watermark (small, 35-45% of sidebar width).
    arrow_w = sidebar_w * 0.4
    aw = slide.shapes.add_picture(
        str(assets.ARROW_UP_PNG()),
        Inches(sidebar_w - arrow_w),
        Inches(0.9),
        width=Inches(arrow_w),
        height=Inches(arrow_w),
    )
    # ~10% opacity per DESIGN.md section 4.
    from lxml import etree
    from pptx.oxml.ns import qn
    blip_fill = aw._element.find(qn("p:blipFill"))
    if blip_fill is not None:
        blip = blip_fill.find(qn("a:blip"))
        if blip is not None:
            etree.SubElement(blip, qn("a:alphaModFix"), {"amt": "10000"})

    # Duration pill at top of sidebar.
    duration = spec.get("duration", "1-3 weeks")
    add_outlined_pill(
        slide,
        text=duration,
        left_in=0.4,
        top_in=1.2,
        height_in=0.36,
    )

    # Two-line service name (thin over heavy).
    service_top = spec.get("service_top", "")
    service_bottom = spec.get("service_bottom", "")
    sn = add_textbox(slide, 0.4, 1.8, sidebar_w - 0.8, 1.8, anchor="top")
    sn.text_frame.word_wrap = True
    p1 = sn.text_frame.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    set_run_typography(
        r1, service_top.upper(), TypeRole(40, 300, 1.0), color_hex=WHITE, weight_override=300
    )
    p2 = sn.text_frame.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    set_run_typography(
        r2, service_bottom.upper(), TypeRole(56, 900, 1.0), color_hex=WHITE, weight_override=900
    )

    # Chat-bubble callout: white rounded card with triangular tail pointing
    # UP at the service-name title above. reference.md component spec.
    chat_text = spec.get("chat_bubble")
    if chat_text:
        add_chat_bubble(
            slide,
            text=chat_text,
            left_in=0.4,
            top_in=4.05,
            width_in=sidebar_w - 0.8,
            height_in=0.7,
            text_color_hex=INK,
            tail_position="top",
            tail_offset_in=0.6,
        )

    # Tagline at bottom of sidebar (lifted to leave clearance for safe area).
    tagline = spec.get("tagline")
    if tagline:
        tb = add_textbox(slide, 0.4, SLIDE_HEIGHT_IN - 1.6, sidebar_w - 0.8, 1.2, anchor="bottom")
        tb.text_frame.word_wrap = True
        tp = tb.text_frame.paragraphs[0]
        tp.alignment = PP_ALIGN.LEFT
        tr = tp.add_run()
        set_run_typography(tr, tagline, TYPE_SCALE["body"], color_hex=WHITE)

    # Main area: WHAT'S INCLUDED card with bullet list.
    main_left = sidebar_w + 0.4
    main_w = SLIDE_WIDTH_IN - main_left - 0.4

    # Product mockup or placeholder.
    add_image_or_placeholder(
        slide,
        image_path=spec.get("image"),
        placeholder_label="[PRODUCT MOCKUP / SCREENSHOTS]",
        left_in=main_left,
        top_in=1.2,
        width_in=main_w,
        height_in=3.2,
        on_dark_surface=True,
    )

    # WHAT'S INCLUDED card (60% width of main area, center-bottom).
    incl_w = main_w * 0.95
    incl_left = main_left + (main_w - incl_w) / 2
    incl_top = 4.7
    incl_h = 2.4
    add_card(
        slide,
        left_in=incl_left,
        top_in=incl_top,
        width_in=incl_w,
        height_in=incl_h,
        fill_color_hex=WHITE,
        radius_pt=14,
    )
    lab = add_textbox(slide, incl_left + 0.3, incl_top + 0.18, incl_w - 0.6, 0.3)
    lp = lab.text_frame.paragraphs[0]
    lp.alignment = PP_ALIGN.LEFT
    lr = lp.add_run()
    set_run_typography(
        lr,
        "WHAT'S INCLUDED",
        TYPE_SCALE["secondary_body"],
        color_hex=CORE_BLUE,
        weight_override=700,
    )
    rpr = lr._r.get_or_add_rPr()
    rpr.set("spc", "100")

    items = spec.get("included")
    add_arrow_bullet_list(
        slide,
        items=items,
        left_in=incl_left + 0.3,
        top_in=incl_top + 0.55,
        width_in=incl_w - 0.6,
        item_height_in=(incl_h - 0.7) / max(1, len(items)),
        text_color_hex=INK,
    )

    # Optional service URL bottom-right (well inside the safe area).
    url = spec.get("url")
    if url:
        url_w = 6.0
        ub = add_textbox(
            slide,
            SLIDE_WIDTH_IN - url_w - 0.7,
            SLIDE_HEIGHT_IN - 0.75,
            url_w,
            0.4,
        )
        ub.text_frame.word_wrap = False
        up = ub.text_frame.paragraphs[0]
        up.alignment = PP_ALIGN.RIGHT
        ur = up.add_run()
        set_run_typography(ur, url, TYPE_SCALE["footer"], color_hex=WHITE)
