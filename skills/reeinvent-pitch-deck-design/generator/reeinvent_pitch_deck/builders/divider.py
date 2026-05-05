"""Section divider. DESIGN.md section 6.2 + reference.md A2.

Full-bleed Signature Gradient. Arrow watermark, section label top-left,
title vertically centered, optional one-line description below.
"""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Pt

from reeinvent_pitch_deck.helpers import (
    add_arrow_watermark,
    add_textbox,
    fill_slide_background_gradient,
    set_run_typography,
)
from reeinvent_pitch_deck.spec import SlideSpec
from reeinvent_pitch_deck.theme import (
    SIGNATURE_GRADIENT,
    SLIDE_HEIGHT_IN,
    SLIDE_WIDTH_IN,
    TYPE_SCALE,
    WHITE,
)


def build(slide: Slide, spec: SlideSpec) -> None:
    fill_slide_background_gradient(slide, SIGNATURE_GRADIENT)
    add_arrow_watermark(slide, on_dark_surface=True, size_in=4.5)

    section_label = spec.get("section_label")
    if section_label:
        eb = add_textbox(slide, 0.8, 0.8, 8.0, 0.5)
        p = eb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        set_run_typography(
            run, section_label.upper(), TYPE_SCALE["eyebrow"], color_hex=WHITE
        )
        rpr = run._r.get_or_add_rPr()
        rpr.set("spc", "150")

    title = spec.get("title")
    title_block = add_textbox(
        slide,
        left_in=0.8,
        top_in=2.0,
        width_in=SLIDE_WIDTH_IN - 1.6,
        height_in=3.5,
        anchor="middle",
    )
    title_block.text_frame.word_wrap = True
    p = title_block.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    set_run_typography(run, title, TYPE_SCALE["section_divider"], color_hex=WHITE)

    description = spec.get("description")
    if description:
        sp = title_block.text_frame.add_paragraph()
        sp.alignment = PP_ALIGN.LEFT
        sp.space_before = Pt(20)
        r = sp.add_run()
        set_run_typography(r, description, TYPE_SCALE["body"], color_hex=WHITE)
