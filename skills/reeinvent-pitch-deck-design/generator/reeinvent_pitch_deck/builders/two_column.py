"""Two-column 50/50 split content slide. DESIGN.md section 6.4.

Off-White. Left column: eyebrow + title + body. Right column: a card or block
of body copy. Both columns left-aligned within their column.
"""

from __future__ import annotations

from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Pt

from reeinvent_pitch_deck.helpers import (
    add_brand_stamp,
    add_card,
    add_gradient_stripe,
    add_textbox,
    fill_run_with_gradient,
    fill_slide_background,
    set_run_typography,
    set_card_text_normautofit,
)
from reeinvent_pitch_deck.spec import SlideSpec
from reeinvent_pitch_deck.theme import (
    INK,
    OFF_WHITE,
    SIGNATURE_GRADIENT,
    SLIDE_HEIGHT_IN,
    SLIDE_WIDTH_IN,
    TYPE_SCALE,
    WHITE,
)


def build(slide: Slide, spec: SlideSpec) -> None:
    fill_slide_background(slide, OFF_WHITE)
    add_brand_stamp(slide, on_dark_surface=False)

    eyebrow = spec.get("eyebrow")
    title = spec.get("title")
    left_body = spec.get("left_body")
    right_body = spec.get("right_body")
    highlight_words = spec.get("highlight_words", [])

    # Eyebrow + stripe.
    if eyebrow:
        eb = add_textbox(slide, 0.6, 0.6, 5.5, 0.4)
        p = eb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        set_run_typography(run, eyebrow.upper(), TYPE_SCALE["eyebrow"], color_hex=INK)
        fill_run_with_gradient(run, SIGNATURE_GRADIENT)
        rpr = run._r.get_or_add_rPr()
        rpr.set("spc", "150")
        approx_w = max(0.6, len(eyebrow) * TYPE_SCALE["eyebrow"].size_pt * 0.011)
        add_gradient_stripe(
            slide, left_in=0.6, top_in=0.92, width_in=approx_w
        )

    # Title (left column, ~5.5 in wide).
    tb = add_textbox(slide, 0.6, 1.4, 5.5, 1.2, anchor="top")
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    role = TYPE_SCALE["title"]
    role40 = type(role)(40, role.weight, role.line_height)
    _emit_title_runs(p, title, highlight_words, INK, role40)

    # Left body.
    lb = add_textbox(slide, 0.6, 2.8, 5.5, 4.0, anchor="top")
    lb.text_frame.word_wrap = True
    lp = lb.text_frame.paragraphs[0]
    lp.alignment = PP_ALIGN.LEFT
    lr = lp.add_run()
    set_run_typography(lr, left_body, TYPE_SCALE["body"], color_hex=INK)

    # Right column: white card sized to hug the body text + symmetric padding.
    # Estimate height from char count and target line width; cap at slide.
    body_chars = len(right_body)
    avg_chars_per_line = 38  # approximate at 20pt body inside ~5 in card width
    lines = max(2, body_chars // avg_chars_per_line + 1)
    body_pt = TYPE_SCALE["body"].size_pt
    line_height = body_pt * TYPE_SCALE["body"].line_height
    body_height_in = (lines * line_height) / 72
    pad_in = 0.45
    card_height = min(SLIDE_HEIGHT_IN - 2.2, body_height_in + pad_in * 2)
    card_top = 2.6  # anchor below the title
    card = add_card(
        slide,
        left_in=7.0,
        top_in=card_top,
        width_in=5.7,
        height_in=card_height,
        fill_color_hex=WHITE,
        radius_pt=16,
    )
    rb = add_textbox(
        slide,
        7.0 + pad_in,
        card_top + pad_in,
        5.7 - pad_in * 2,
        card_height - pad_in * 2,
        anchor="top",
    )
    rb.text_frame.word_wrap = True
    rbp = rb.text_frame.paragraphs[0]
    rbp.alignment = PP_ALIGN.LEFT
    rbr = rbp.add_run()
    set_run_typography(rbr, right_body, TYPE_SCALE["body"], color_hex=INK)


def _emit_title_runs(paragraph, title: str, highlights: list[str], color: str, role) -> None:
    if not highlights:
        run = paragraph.add_run()
        set_run_typography(run, title, role, color_hex=color)
        return
    remaining = title
    for hl in highlights:
        idx = remaining.find(hl)
        if idx == -1:
            continue
        before = remaining[:idx]
        if before:
            run = paragraph.add_run()
            set_run_typography(run, before, role, color_hex=color)
        run = paragraph.add_run()
        set_run_typography(run, hl, role, color_hex=color)
        fill_run_with_gradient(run, SIGNATURE_GRADIENT)
        remaining = remaining[idx + len(hl):]
    if remaining:
        run = paragraph.add_run()
        set_run_typography(run, remaining, role, color_hex=color)
